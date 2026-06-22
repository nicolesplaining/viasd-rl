#!/usr/bin/env python3
"""Build <deck>.pptx from <deck>.pdf (slide images) + <deck>.tex (\\note speaker notes).
Each PPTX slide = the rendered beamer slide as a full-bleed image, with the speaker note
attached as a native PowerPoint note (editable in the Notes pane). Usage: build_pptx.py <deck>"""
import os
import re
import sys

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

deck = sys.argv[1]
pdf, tex = deck + ".pdf", deck + ".tex"

# ---- extract \note{...} blocks in document order (brace-aware) ----
s = open(tex).read()
notes_raw, i = [], 0
while True:
    j = s.find("\\note{", i)
    if j == -1:
        break
    k, depth = j + 6, 1
    while depth > 0:
        c = s[k]
        depth += (c == "{") - (c == "}")
        k += 1
    notes_raw.append(s[j + 6:k - 1])
    i = k

def clean(t):
    t = t.replace("\\textsuperscript{4}", "⁴")          # S\textsuperscript{4}D -> S⁴D
    for _ in range(4):                                        # unwrap formatting commands
        t = re.sub(r"\\(?:emph|textbf|textit|structure|texttt|mathbf|text|underline)\{([^{}]*)\}", r"\1", t)
    t = t.replace("``", '"').replace("''", '"').replace("---", "—").replace("--", "–")
    t = re.sub(r"\$([^$]*)\$", r"\1", t)                      # strip math delimiters
    t = re.sub(r"\\times", "x", t)
    t = re.sub(r"\\[a-zA-Z]+", "", t)                         # drop leftover commands
    t = t.replace("{", "").replace("}", "")
    return re.sub(r"\s+", " ", t).strip()

notes = [clean(n) for n in notes_raw]

# ---- render PDF pages to PNG ----
doc = fitz.open(pdf)
outdir = f"/tmp/{deck}_png"
os.makedirs(outdir, exist_ok=True)
mat = fitz.Matrix(2.6, 2.6)  # ~190 dpi
imgs = []
for pi in range(doc.page_count):
    p = f"{outdir}/p{pi:02d}.png"
    doc[pi].get_pixmap(matrix=mat).save(p)
    imgs.append(p)

assert len(imgs) == len(notes), f"page/note mismatch: {len(imgs)} pages vs {len(notes)} notes"

# ---- assemble 16:9 PPTX ----
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]
for img, note in zip(imgs, notes):
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
    slide.notes_slide.notes_text_frame.text = note
prs.save(deck + ".pptx")
print(f"{deck}.pptx: {len(imgs)} slides, notes attached ({sum(len(n.split()) for n in notes)} note words)")
